import sys

# Шаг 0: Проверка зависимостей
missing_libraries = []
try:
    import fitz  # PyMuPDF
except ImportError:
    missing_libraries.append("pymupdf")

try:
    import pptx
except ImportError:
    missing_libraries.append("python-pptx")

if missing_libraries:
    print("❌ Ошибка: В системе отсутствуют необходимые библиотеки Python.")
    print("📋 Для их установки выполните команду в терминале:")
    print(f"\n    pip install {' '.join(missing_libraries)}\n")
    sys.exit(1)

import os
import argparse
import subprocess
import shutil
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor

def make_dark_mode(pptx_path, temp_output_path):
    """Создает копию презентации с черным фоном, белым текстом и без фоновых картинок."""
    prs = Presentation(pptx_path)
    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    
    for slide in prs.slides:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BLACK
        
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                slide_area = prs.slide_width * prs.slide_height
                shape_area = shape.width * shape.height
                if shape_area / slide_area > 0.8:
                    sp = shape._element
                    sp.getparent().remove(sp)
                    continue

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = WHITE

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = WHITE
                                
    prs.save(temp_output_path)

def get_resolution_multiplier(quality_str, page):
    rect = page.rect
    long_side = max(rect.width, rect.height)
    
    if quality_str == "2k":
        target = 2560
    elif quality_str == "4k":
        target = 3840
    else:
        return 2.0
        
    return target / long_side

def pptx_to_pdf_crossplatform(pptx_path, output_dir):
    """Автоматически определяет ОС и конвертирует PPTX в PDF через LibreOffice."""
    mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    
    if os.path.exists(mac_path):
        libreoffice_path = mac_path
    elif shutil.which("soffice") is not None:
        libreoffice_path = "soffice"
    else:
        raise FileNotFoundError(
            "LibreOffice не найден! На Linux установите через: sudo apt install libreoffice"
        )

    cmd = [
        libreoffice_path,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(pptx_path)
    ]
    subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
    return output_dir / f"{pptx_path.stem}.pdf"

def pdf_to_png_fast(pdf_path, output_dir, quality):
    doc = fitz.open(pdf_path)
    total_pages = len(doc)
    z_fill_len = max(2, len(str(total_pages)))
    
    created_files = []
    for page_num in range(total_pages):
        page = doc.load_page(page_num)
        zoom = get_resolution_multiplier(quality, page)
        mat = fitz.Matrix(zoom, zoom)
        
        pix = page.get_pixmap(matrix=mat)
        slide_index = str(page_num + 1).zfill(z_fill_len)
        png_path = output_dir / f"slide_{slide_index}.png"
        pix.save(str(png_path))
        created_files.append(png_path)
        
    doc.close()
    return total_pages, created_files

def create_zip_archive(files_to_zip, zip_output_path):
    """Создает ZIP-архив и упаковывает туда файлы."""
    with zipfile.ZipFile(zip_output_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
        for file in files_to_zip:
            zipf.write(file, arcname=file.name)

def process_file(pptx_path, args):
    print(f"\n🚀 Обработка файла: {pptx_path.name}")
    
    base_dir = Path(args.output_dir or pptx_path.parent)
    file_output_dir = base_dir / f"{pptx_path.stem}_output"
    file_output_dir.mkdir(parents=True, exist_ok=True)
    
    current_pptx = pptx_path
    temp_dark_pptx = None
    zip_created_successfully = False

    try:
        if args.dark_mode:
            print(" 🎨 Применение темной темы (черный ...)")
            temp_dark_pptx = file_output_dir / f"temp_dark_{pptx_path.name}"
            make_dark_mode(pptx_path, temp_dark_pptx)
            current_pptx = temp_dark_pptx

        # 1. Конвертация в PDF
        pdf_path = pptx_to_pdf_crossplatform(current_pptx, file_output_dir)
        print(" 📄 Конвертация в PDF успешна.")
        
        # 2. Конвертация в PNG
        print(f" 🖼️ Конвертация в PNG (Качество: {args.quality})...")
        total_slides, generated_pngs = pdf_to_png_fast(pdf_path, file_output_dir, args.quality)
        print(f" ✅ Успешно создано слайдов: {total_slides}")
        
        # 3. Архивирование в ZIP
        if args.zip:
            zip_name = f"{pptx_path.stem}_output.zip"
            zip_path = base_dir / zip_name
            print(f" 📦 Архивирование слайдов в {zip_name}...")
            create_zip_archive(generated_pngs, zip_path)
            print(" 📦 Создание архива завершено.")
            zip_created_successfully = True
        
        # Удаление временного PDF
        if not args.keep_pdf and pdf_path.exists():
            pdf_path.unlink()
            print(" 🗑️ Временный PDF удален.")
            
        # Удаление временного инвертированного PPTX
        if temp_dark_pptx and temp_dark_pptx.exists():
            temp_dark_pptx.unlink()

        # 4. Очистка исходной папки с картинками, если взведен флаг --clean
        if args.clean:
            if zip_created_successfully:
                print(f" 🧹 Удаление папки с несжатыми PNG для экономии места...")
                shutil.rmtree(file_output_dir)
                print(" 🧹 Папка успешно удалена. Оставлен только чистый ZIP-архив.")
            else:
                print(" ⚠️ Предупреждение: Ключ --clean работает только вместе с ключом --zip! Папка не удалена.")

    except Exception as e:
        print(f" ❌ Ошибка при обработке {pptx_path.name}: {e}")

def main():
    parser = argparse.ArgumentParser(description="Универсальный конвертер PPTX -> PDF -> PNG.")
    
    parser.add_argument("-f", "--file", type=str, help="Путь к конкретному файлу PPTX.")
    parser.add_argument("-d", "--dir", type=str, default=".", help="Путь к папке с презентациями.")
    parser.add_argument("-o", "--output-dir", type=str, help="Папка для сохранения результатов.")
    parser.add_argument("-a", "--all", action="store_true", help="Обработать ВСЕ файлы pptx в папке.")
    parser.add_argument("--keep-pdf", action="store_true", help="Сохранить промежуточный PDF файл.")
    parser.add_argument("-q", "--quality", choices=["standard", "2k", "4k"], default="standard", help="Качество выходных PNG.")
    parser.add_argument("--dark-mode", action="store_true", help="Включить темную тему.")
    parser.add_argument("--zip", action="store_true", help="Упаковать все готовые PNG картинки в ZIP-архив.")
    
    # Новый ключ для автоматической очистки папки картинок
    parser.add_argument("--clean", action="store_true", help="Удалить папку с исходными PNG после успешного создания ZIP-архива.")

    args = parser.parse_args()
    files_to_process = []
    
    if args.file:
        specific_file = Path(args.file).resolve()
        if specific_file.exists() and specific_file.suffix.lower() in ['.pptx', '.ppt']:
            files_to_process.append(specific_file)
        else:
            print(f"❌ Ошибка: Файл '{args.file}' не найден.")
            return
    else:
        target_dir = Path(args.dir).resolve()
        if not target_dir.exists():
            print(f"❌ Ошибка: Папка '{args.dir}' не существует.")
            return
            
        all_pptx = [f for f in target_dir.glob("*.[pP][pP][tT]*") if not f.name.startswith("~$")]
        
        if not all_pptx:
            print(f"❌ Ошибка: В папке '{target_dir}' не найдено презентаций.")
            return
            
        if args.all:
            files_to_process = all_pptx
        else:
            files_to_process = [all_pptx[0]]

    for pptx_file in files_to_process:
        process_file(pptx_file, args)
        
    print("\n🎉 Работа скрипта завершена!")

if __name__ == "__main__":
    main()
