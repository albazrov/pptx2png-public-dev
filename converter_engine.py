import os
import asyncio
import subprocess
import shutil
import zipfile
import re
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
import fitz  # PyMuPDF

class FakeArgs:
    def __init__(self, quality, keep_pdf, dark_mode=True, zip_mode=True, clean=True, output_dir=None):
        self.quality = quality
        self.keep_pdf = keep_pdf
        self.dark_mode = dark_mode
        self.zip = zip_mode
        self.clean = clean
        self.output_dir = output_dir

def make_dark_mode(pptx_path, temp_output_path):
    """Создает копию презентации с черным фоном, белым текстом и без фоновых картинок."""
    prs = Presentation(pptx_path)
    BLACK = RGBColor(0, 0, 0)
    WHITE = RGBColor(255, 255, 255)
    
    for slide in prs.slides:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BLACK
        
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                slide_area = prs.slide_width * prs.slide_height
                shape_area = shape.width * shape.height
                if shape_area / slide_area > 0.8:
                    sp = shape._element
                    sp.getparent().remove(sp)
                    continue

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = WHITE

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = WHITE
                                
    prs.save(temp_output_path)

def get_resolution_multiplier(quality_str, page):
    rect = page.rect
    long_side = max(rect.width, rect.height)
    if quality_str == "2k":
        target = 2560
    elif quality_str == "4k":
        target = 3840
    else:
        return 2.0
    return target / long_side

def pptx_to_pdf_crossplatform(pptx_path, output_dir):
    mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.exists(mac_path):
        libreoffice_path = mac_path
    elif shutil.which("soffice") is not None:
        libreoffice_path = "soffice"
    else:
        raise FileNotFoundError("LibreOffice не найден в системе!")

    cmd = [
        libreoffice_path, "--headless",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(pptx_path)
    ]
    subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
    return output_dir / f"{pptx_path.stem}.pdf"

def pdf_to_png_fast(pdf_path, output_dir, quality):
    doc = fitz.open(pdf_path)
    total_pages = len(doc)
    z_fill_len = max(2, len(str(total_pages)))
    
    created_files = []
    for page_num in range(total_pages):
        page = doc.load_page(page_num)
        zoom = get_resolution_multiplier(quality, page)
        mat = fitz.Matrix(zoom, zoom)
        
        pix = page.get_pixmap(matrix=mat)
        slide_index = str(page_num + 1).zfill(z_fill_len)
        png_path = output_dir / f"slide_{slide_index}.png"
        pix.save(str(png_path))
        created_files.append(png_path)
        
    doc.close()
    return total_pages, created_files

def extract_zip_if_needed(zip_path, extract_dir):
    """Распаковывает ZIP архив и ищет в нём файл презентации."""
    with zipfile.ZipFile(zip_path, 'r') as zip_ref:
        zip_ref.extractall(extract_dir)
    for ext in ("*.pptx", "*.PPTX", "*.ppt", "*.PPT"):
        for file in extract_dir.glob(ext):
            if not file.name.startswith("~$"):
                return file
    return None

def convert_to_direct_download(url: str) -> str:
    """Преобразует ссылки Google Drive/Slides в точную ссылку на скачивание."""
    url = url.strip()
    
    # 1. Ссылки на веб-презентации Google Slides (://google.com...)
    if "docs.google.com/presentation" in url:
        match = re.search(r'/presentation/d/([a-zA-Z0-9_-]+)', url)
        if match:
            file_id = match.group(1)
            return f"https://docs.google.com/presentation/d/{file_id}/export/pptx"

    # 2. Стандартные файлы презентаций на Google Диске (://google.com...)
    if "://google.com" in url:
        match = re.search(r'/file/d/([a-zA-Z0-9_-]+)', url)
        if match:
            file_id = match.group(1)
            return f"https://google.com{file_id}"
            
    # 3. Ссылки с прямым параметром ?id=... или &id=...
    match_id = re.search(r'[?&]id=([a-zA-Z0-9_-]+)', url)
    if match_id:
        file_id = match_id.group(1)
        return f"https://google.com{file_id}"
        
    return url

def process_file_local(pptx_path, args):
    """
    Синхронное ядро обработки (бывший метод из convert.py).
    Возвращает путь к созданному ZIP-архиву.
    """
    file_output_dir = Path(args.output_dir) / f"{pptx_path.stem}_output"
    file_output_dir.mkdir(parents=True, exist_ok=True)
    
    current_pptx = pptx_path
    temp_dark_pptx = None
    zip_path = None
    
    try:
        if args.dark_mode:
            temp_dark_pptx = file_output_dir / f"temp_dark_{pptx_path.name}"
            make_dark_mode(pptx_path, temp_dark_pptx)
            current_pptx = temp_dark_pptx

        pdf_path = pptx_to_pdf_crossplatform(current_pptx, file_output_dir)
        total_slides, generated_pngs = pdf_to_png_fast(pdf_path, file_output_dir, args.quality)
        
        if args.zip:
            # Создаём ZIP архив с строго определённым именем
            zip_path = Path(args.output_dir) / f"{pptx_path.stem}_output.zip"
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for file in generated_pngs:
                    zipf.write(file, arcname=file.name)
        
        if not args.keep_pdf and pdf_path.exists():
            pdf_path.unlink()
        if temp_dark_pptx and temp_dark_pptx.exists():
            temp_dark_pptx.unlink()
        
        # Возвращаем путь к созданному ZIP-архиву
        return zip_path
            
    except Exception as e:
        # Если произошла ошибка, чистим временные файлы
        if temp_dark_pptx and temp_dark_pptx.exists():
            temp_dark_pptx.unlink()
        raise e
