import aiohttp
import logging
import shutil
import asyncio  # Добавлен для перехвата TimeoutError
from pathlib import Path
from typing import Tuple, List, Optional
from aiogram import Bot, types
from pptx import Presentation

# Импортируем ваш существующий движок рендеринга
import converter_engine

# ==========================================
# 1. ПРОВЕРКА ОРФОГРАФИИ (С РАЗДЕЛЕНИЕМ СТАТУСОВ)
# ==========================================

async def check_spelling(text_list: list) -> Tuple[bool, str]:
    """
    Асинхронно проверяет список текстов слайдов через Яндекс.Спеллер API.
    
    :param text_list: Список строк с текстом каждого слайда
    :return: Кортеж (успех_проверки, текст_отчёта)
             - (True, report) - проверка выполнена, отчёт с ошибками или "ошибок нет"
             - (False, error_message) - проверка не удалась, текст ошибки
    """
    # Если список пуст или содержит только пустые строки
    if not text_list or all(not text.strip() for text in text_list):
        return True, "ℹ️ **В презентации не найден текст для проверки.**"
    
    url = "https://speller.yandex.net/services/spellservice.json/checkText"
    report_lines = []
    check_success = False
    has_errors = False
    
    async with aiohttp.ClientSession() as session:
        for idx, slide_text in enumerate(text_list, start=1):
            if not slide_text.strip():
                continue
                
            payload = {
                "text": slide_text,
                "options": 518  # Игнорировать URL, email и римские цифры
            }
            
            try:
                async with session.post(url, data=payload, timeout=10) as response:
                    if response.status == 200:
                        results = await response.json()
                        check_success = True
                        if results:
                            has_errors = True
                            report_lines.append(f"📋 **Слайд №{idx}:**")
                            for error in results:
                                word = error.get("word", "неизвестное слово")
                                s_list = error.get("s", [])
                                suggestions = f" ➔ возможно: `{', '.join(s_list)}`" if s_list else " (нет вариантов)"
                                report_lines.append(f"  • Опечатка в `{word}`{suggestions}")
                    else:
                        # Не 200 ответ — возвращаем ошибку
                        return False, f"❌ **Ошибка сервиса проверки орфографии:**\nКод ответа: {response.status}\nПопробуйте позже или используйте конвертацию без проверки."
                        
            except asyncio.TimeoutError:
                # Перехватываем asyncio.TimeoutError (включает aiohttp серверные таймауты)
                return False, "❌ **Таймаут при проверке орфографии.**\nСервис не отвечает. Попробуйте позже."
            except aiohttp.ClientError as e:
                # Перехватываем все клиентские ошибки aiohttp (ConnectionError, ClientResponseError и т.д.)
                return False, f"❌ **Ошибка сети при проверке орфографии:**\n{str(e)}\nПопробуйте позже."
            except Exception as e:
                # Перехватываем любые другие неожиданные ошибки
                logging.error(f"Неизвестная ошибка Яндекс.Спеллера на слайде {idx}: {e}", exc_info=True)
                return False, f"❌ **Неизвестная ошибка при проверке:**\n{str(e)}\nПопробуйте позже."
    
    # Если ни один слайд не был отправлен на проверку
    if not check_success:
        return False, "❌ **Не удалось проверить текст.**\nСервис не ответил ни на один запрос."
    
    # Проверка выполнена успешно
    if has_errors:
        header = "⚠️ **Внимание! На слайдах обнаружены возможные опечатки:**\n\n"
        return True, header + "\n".join(report_lines)
    else:
        return True, "✨ **Орфографических ошибок не найдено!** Текст чистый."


# ==========================================
# 2. ИЗВЛЕЧЕНИЕ ТЕКСТА ИЗ PPTX (С ОТДЕЛЬНЫМ СТАТУСОМ)
# ==========================================

def extract_text_from_pptx(file_path: str) -> Tuple[bool, List[str]]:
    """
    Извлекает весь текст из блоков презентации, группируя его по слайдам.
    
    :param file_path: Путь к локальному файлу .pptx
    :return: Кортеж (успех_извлечения, список_текстов_слайдов)
             - (True, list) - текст успешно извлечён (может быть пустым)
             - (False, []) - ошибка при извлечении
    """
    try:
        prs = Presentation(file_path)
        presentation_text = []
        
        for slide in prs.slides:
            slide_pieces = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_pieces.append(shape.text)
            presentation_text.append(" ".join(slide_pieces))
        
        return True, presentation_text
        
    except Exception as e:
        logging.error(f"Ошибка извлечения текста из PPTX через python-pptx: {e}")
        return False, []


# ==========================================
# 3. СКАЧИВАНИЕ ПО ССЫЛКЕ
# ==========================================

async def download_file_by_url(url: str, destination: Path, status_message: types.Message) -> bool:
    """ Скачивает презентацию по HTTP-ссылке напрямую в RAM-диск (SHM). """
    try:
        async with aiohttp.ClientSession() as session:
            async with session.get(url, timeout=30) as response:
                if response.status != 200:
                    logging.error(f"Ошибка скачивания по ссылке. Статус: {response.status}")
                    return False
                
                with open(destination, "wb") as f:
                    f.write(await response.read())
                return True
    except asyncio.TimeoutError:
        logging.error(f"Таймаут при скачивании по ссылке: {url}")
        return False
    except aiohttp.ClientError as e:
        logging.error(f"Ошибка клиента при скачивании: {e}")
        return False
    except Exception as e:
        logging.error(f"Исключение при скачивании по URL: {e}")
        return False

# ==========================================
# 4. ОСНОВНОЙ КОНВЕЙЕР ОБРАБОТКИ
# ==========================================

async def core_pipeline(downloaded_file_path: Path, status_message: types.Message, user_id: int, user_mgr):
    """
    Основной конвейер конвертации презентации.
    Возвращает кортеж (путь_к_zip, путь_к_pdf) или (None, None) при ошибке.
    """
    work_dir = downloaded_file_path.parent
    is_zip = downloaded_file_path.suffix.lower() == '.zip'
    cfg = user_mgr.get_user_config(user_id)

    try:
        # 1. Если загружен ZIP — распаковываем и находим PPTX
        if is_zip:
            await status_message.edit_text("📦 Распаковка ZIP...")
            pptx_path = converter_engine.extract_zip_if_needed(downloaded_file_path, work_dir)
            if not pptx_path:
                await status_message.edit_text("❌ Внутри ZIP не найдено .pptx.")
                return None, None
        else:
            pptx_path = downloaded_file_path

        # 2. Запускаем конвертацию
        await status_message.edit_text(f"⏳ Конвертация через LibreOffice в RAM...\n(Качество: {cfg['quality'].upper()})")
        
        clean_folder = not cfg["keep_pdf"]
        args = converter_engine.FakeArgs(
            quality=cfg["quality"], 
            keep_pdf=cfg["keep_pdf"], 
            dark_mode=True, 
            zip_mode=True, 
            clean=clean_folder, 
            output_dir=str(work_dir)
        )
        
        # Вызов движка в фоновом пуле потоков
        # ВАЖНО: process_file_local теперь возвращает путь к созданному ZIP
        expected_zip = await asyncio.to_thread(converter_engine.process_file_local, pptx_path, args)
        
        # Проверяем, что ZIP действительно создан
        if not expected_zip or not expected_zip.exists():
            await status_message.edit_text("❌ Ошибка: ZIP архив не создан.")
            return None, None
        
        # 3. Обработка PDF (если пользователь хочет его сохранить)
        final_pdf_path = None
        if cfg["keep_pdf"]:
            png_folder = next((d for d in work_dir.iterdir() if d.is_dir() and d.name.endswith("_output")), None)
            if png_folder:
                pdf_file = next(png_folder.glob("*.pdf"), None)
                if pdf_file:
                    final_pdf_path = work_dir / f"{pptx_path.stem}.pdf"
                    shutil.move(str(pdf_file), str(final_pdf_path))
                # Удаляем временную папку с PNG и PDF
                shutil.rmtree(png_folder)

        # 4. Если это был ZIP, удаляем исходный загруженный архив
        if is_zip and downloaded_file_path.exists():
            downloaded_file_path.unlink()

        return expected_zip, final_pdf_path
        
    except Exception as e:
        logging.error(f"Критическая ошибка ядра: {e}", exc_info=True)
        return None, None
